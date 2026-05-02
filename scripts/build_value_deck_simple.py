"""Build the Silvers Fun value deck (one-shot, spec-matching version)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "assets" / "screenshots"
BRAND = ROOT / "assets" / "branding"
OUT_PPTX = ROOT / "docs" / "presentations" / "silvers_fun_value_deck.pptx"
OUT_MD = ROOT / "docs" / "presentations" / "silvers_fun_value_deck_summary.md"

CREAM = RGBColor(0xFA, 0xF4, 0xE8)
SOFT_LAV = RGBColor(0xEF, 0xEA, 0xF7)
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
CORAL = RGBColor(0xE2, 0x7D, 0x60)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
INK = RGBColor(0x2C, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=CREAM):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False


def add_accent_bars(slide):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    top.line.fill.background()
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.shadow.inherit = False
    bot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.12), SW, Inches(0.12))
    bot.line.fill.background()
    bot.fill.solid()
    bot.fill.fore_color.rgb = CORAL
    bot.shadow.inherit = False


def add_text(slide, left, top, width, height, text, size=24, bold=False,
             color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, size=24, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_phone(slide, left, top, height_in, image_path):
    aspect = 1080 / 2340
    height = Inches(height_in)
    inner_w = int(height * aspect)
    pad = Inches(0.12)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, inner_w + pad * 2,
        height + pad * 2)
    frame.adjustments[0] = 0.10
    frame.fill.solid()
    frame.fill.fore_color.rgb = NAVY
    frame.line.color.rgb = NAVY
    frame.shadow.inherit = False
    slide.shapes.add_picture(
        str(image_path), left + pad, top + pad, width=inner_w, height=height)
    return inner_w + pad * 2


def add_phone_placeholder(slide, left, top, height_in, label):
    """Phone-shaped placeholder for a screenshot to be added later."""
    aspect = 1080 / 2340
    height = Inches(height_in)
    inner_w = int(height * aspect)
    pad = Inches(0.12)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
        inner_w + pad * 2, height + pad * 2)
    frame.adjustments[0] = 0.10
    frame.fill.solid()
    frame.fill.fore_color.rgb = NAVY
    frame.line.color.rgb = NAVY
    frame.shadow.inherit = False
    inner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + pad, top + pad, inner_w, height)
    inner.adjustments[0] = 0.08
    inner.fill.solid()
    inner.fill.fore_color.rgb = SOFT_LAV
    inner.line.color.rgb = CORAL
    inner.line.width = Pt(1.5)
    inner.shadow.inherit = False
    tf = inner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Add screenshot"
    r.font.name = "Calibri"
    r.font.size = Pt(13)
    r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Calibri"
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return inner_w + pad * 2


def add_rounded_image(slide, left, top, width, height, image_path):
    pad = Inches(0.1)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width + pad * 2,
        height + pad * 2)
    frame.adjustments[0] = 0.06
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = NAVY
    frame.line.width = Pt(1.5)
    frame.shadow.inherit = False
    slide.shapes.add_picture(
        str(image_path), left + pad, top + pad, width=width, height=height)


def add_header(slide, eyebrow, title, eyebrow_color=CORAL):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.45),
             eyebrow.upper(), size=16, bold=True, color=eyebrow_color)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.95),
             title, size=40, bold=True, color=NAVY)


def add_page_number(slide, n, total=10):
    add_text(slide, SW - Inches(1.4), SH - Inches(0.55), Inches(1.2),
             Inches(0.35), f"{n} / {total}", size=12, color=NAVY,
             align=PP_ALIGN.RIGHT)


# Slide 1: Title
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
panel = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.0),
    Inches(5.7), Inches(5.5))
panel.adjustments[0] = 0.05
panel.fill.solid()
panel.fill.fore_color.rgb = SOFT_LAV
panel.line.fill.background()
panel.shadow.inherit = False
fg_w = Inches(5.0)
fg_h = Inches(5.0 * 500 / 1024)
add_rounded_image(s, Inches(7.35), Inches(2.6),
                  fg_w, fg_h, BRAND / "feature_graphic_1024x500.png")
s.shapes.add_picture(
    str(BRAND / "app_icon_play_store_512.png"),
    Inches(0.8), Inches(1.0), width=Inches(1.4), height=Inches(1.4))
add_text(s, Inches(0.8), Inches(2.6), Inches(6.0), Inches(1.5),
         "Silvers Fun", size=72, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(4.0), Inches(6.0), Inches(2.0),
         "Helping seniors connect, meet,\nand enjoy life together.",
         size=28, color=INK)
add_text(s, Inches(0.8), Inches(5.7), Inches(6), Inches(0.4),
         "A warm social companion app for the 65+ community",
         size=16, color=TEAL, bold=True)
add_page_number(s, 1)

# Slide 2: The problem
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "The problem", "Connection gets harder with age")
add_bullets(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.5), [
    "Many seniors want companionship, conversation, and shared activities.",
    "Existing apps feel built for younger users — confusing, fast, or unsafe.",
    "Moving, retiring, or losing a partner can shrink social circles.",
    "Families want a safe, comfortable way for older relatives to make friends.",
], size=26)
add_page_number(s, 2)

# Slide 3: Who it is for
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Who it is for", "Built around the 65+ community")
groups = [
    ("Seniors 65+",
     "Looking for friends, companions, and group activities.", CORAL),
    ("Families",
     "Helping a parent or grandparent stay socially active.", TEAL),
    ("Senior communities",
     "Centers and residences that host meetups and events.", NAVY),
    ("Hobby & interest groups",
     "Walking, cooking, music, faith, language exchange.", CORAL),
]
card_w = Inches(2.95)
card_h = Inches(3.6)
gap = Inches(0.15)
start_x = Inches(0.55)
top = Inches(2.3)
for i, (head, body, accent) in enumerate(groups):
    x = start_x + (card_w + gap) * i
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, card_h)
    card.adjustments[0] = 0.07
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(2)
    card.shadow.inherit = False
    stripe = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, top, card_w, Inches(0.25))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    stripe.shadow.inherit = False
    add_text(s, x + Inches(0.2), top + Inches(0.5), card_w - Inches(0.4),
             Inches(0.8), head, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), top + Inches(1.4), card_w - Inches(0.4),
             Inches(2.0), body, size=18, color=INK)
add_page_number(s, 3)

# Slide 4: Value & features (merged: value statement + 6-feature grid)
s = prs.slides.add_slide(BLANK)
add_bg(s, SOFT_LAV)
add_accent_bars(s)
add_header(s, "What it does", "Real connections, made simple")
add_text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.5),
         "Silvers Fun helps seniors discover people with shared interests, "
         "join meetups, and build meaningful connections.",
         size=20, color=INK)
features = [
    ("Profiles", "Photos and interests that show who you are.", CORAL),
    ("Discover", "Browse people who share your interests.", TEAL),
    ("Meetups", "Create or join group activities nearby.", NAVY),
    ("Likes", "Friendly mutual interest — no swiping.", CORAL),
    ("Chat", "Simple messaging once both sides like each other.", TEAL),
    ("English & 한국어", "Full English and Korean support.", NAVY),
]
fcol_w = Inches(3.94)
fcol_h = Inches(2.05)
fcol_gap_x = Inches(0.20)
fcol_gap_y = Inches(0.20)
fcol_start_x = Inches(0.55)
fcol_start_y = Inches(2.55)
for i, (head, body, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    fx = fcol_start_x + (fcol_w + fcol_gap_x) * col
    fy = fcol_start_y + (fcol_h + fcol_gap_y) * row
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy, fcol_w, fcol_h)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    box.shadow.inherit = False
    badge = s.shapes.add_shape(
        MSO_SHAPE.OVAL, fx + Inches(0.25), fy + Inches(0.32),
        Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    badge.shadow.inherit = False
    add_text(s, fx + Inches(0.25), fy + Inches(0.34), Inches(0.5),
             Inches(0.5), str(i + 1), size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, fx + Inches(0.90), fy + Inches(0.30), fcol_w - Inches(1.1),
             Inches(0.55), head, size=20, bold=True, color=NAVY)
    add_text(s, fx + Inches(0.30), fy + Inches(1.05), fcol_w - Inches(0.5),
             Inches(0.85), body, size=15, color=INK)
add_page_number(s, 4)

# Slide 5: Getting started — easy onboarding with one-tap Google Sign-In
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Getting started", "Easy to start, easy to use")
add_text(s, Inches(0.7), Inches(1.80), Inches(11.9), Inches(0.35),
         "Sign in once and you're ready to meet people — "
         "no usernames or passwords to remember.",
         size=18, color=INK)
# Phone-shaped placeholder for the sign-in screen on the left
add_phone_placeholder(s, Inches(1.6), Inches(2.30), 4.85, "Sign-in screen")
add_text(s, Inches(1.6), Inches(7.20), Inches(2.6), Inches(0.25),
         "One-tap sign-in", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
# Right column: 3 onboarding steps
onboard = [
    ("Sign in with Google",
     "One tap. No passwords or usernames to remember.", CORAL),
    ("Set up your profile",
     "Add a few photos and the interests that matter to you.", TEAL),
    ("Discover and meet",
     "Browse people, join meetups, start friendly conversations.", NAVY),
]
ox = Inches(5.6)
oy = Inches(2.40)
ow = Inches(7.3)
oh = Inches(1.50)
ogap = Inches(0.20)
for i, (head, body, accent) in enumerate(onboard):
    y = oy + (oh + ogap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ox, y, ow, oh)
    box.adjustments[0] = 0.10
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    box.shadow.inherit = False
    badge = s.shapes.add_shape(
        MSO_SHAPE.OVAL, ox + Inches(0.30), y + Inches(0.40),
        Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    badge.shadow.inherit = False
    add_text(s, ox + Inches(0.30), y + Inches(0.42), Inches(0.7),
             Inches(0.7), str(i + 1), size=24, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, ox + Inches(1.15), y + Inches(0.25), ow - Inches(1.35),
             Inches(0.55), head, size=22, bold=True, color=NAVY)
    add_text(s, ox + Inches(1.15), y + Inches(0.80), ow - Inches(1.35),
             Inches(0.65), body, size=15, color=INK)
add_page_number(s, 5)

# Slide 6: Explore people and activities
# (combines Discover + Meetups + Create Meetup)
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Discover & meetups", "Explore people and activities")
add_text(s, Inches(0.7), Inches(1.80), Inches(11.9), Inches(0.35),
         "Discover people with shared interests, browse meetups, "
         "and create or join activities.",
         size=18, color=INK)
phone_h_3 = 4.85
# Match add_phone's internal width math so frame_w_3 stays accurate.
frame_w_3 = int(Inches(phone_h_3) * 1080 / 2340) + Inches(0.24)
gap_3 = Inches(0.45)
total_3 = frame_w_3 * 3 + gap_3 * 2
start_x_3 = (SW - total_3) // 2
y_phone_3 = Inches(2.20)
shots_6 = [
    (SHOTS / "01_discover.jpg", "Discover"),
    (SHOTS / "03_meetups.jpg", "Meetups"),
    (SHOTS / "04_create_meetup.jpg", "Create meetup"),
]
for i, (path, caption) in enumerate(shots_6):
    px = start_x_3 + (frame_w_3 + gap_3) * i
    add_phone(s, px, y_phone_3, phone_h_3, path)
    add_text(s, px, Inches(7.30), frame_w_3, Inches(0.25), caption,
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_page_number(s, 6)

# Slide 7: Personalize your experience
# (profile editing + photos + language + Likes — Likes screenshot pending)
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Account & settings", "Personalize your experience")
add_text(s, Inches(0.7), Inches(1.80), Inches(11.9), Inches(0.35),
         "Set up your profile, photos, language, and how connections start.",
         size=18, color=INK)
phone_h_4 = 4.5
frame_w_4 = int(Inches(phone_h_4) * 1080 / 2340) + Inches(0.24)
gap_4 = Inches(0.30)
total_4 = frame_w_4 * 4 + gap_4 * 3
start_x_4 = (SW - total_4) // 2
y_phone_4 = Inches(2.25)
slide7_items = [
    ("phone", SHOTS / "06_you_profile.jpg", "Profile"),
    ("phone", SHOTS / "08_edit_photos.jpg", "Photos"),
    ("phone", SHOTS / "07_language_picker.jpg", "Language"),
    ("placeholder", "Likes", "Likes"),
]
for i, (kind, asset, caption) in enumerate(slide7_items):
    px = start_x_4 + (frame_w_4 + gap_4) * i
    if kind == "phone":
        add_phone(s, px, y_phone_4, phone_h_4, asset)
    else:
        add_phone_placeholder(s, px, y_phone_4, phone_h_4, asset)
    add_text(s, px, Inches(7.10), frame_w_4, Inches(0.25), caption,
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_page_number(s, 7)

# Slide 8: Connect comfortably and safely
# (Chat + Like → Mutual → Chat flow)
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Friendly chat", "Connect comfortably and safely")
add_text(s, Inches(0.7), Inches(1.80), Inches(11.9), Inches(0.35),
         "Chat opens only after both people like each other — "
         "no unwanted messages.",
         size=18, color=INK)
# Large chat phone on the left
add_phone(s, Inches(1.6), Inches(2.30), 4.85, SHOTS / "05_chat.jpg")
add_text(s, Inches(1.6), Inches(7.20), Inches(2.6), Inches(0.25),
         "Friendly chat", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
# Right column: Like → Mutual → Chat flow
flow_steps = [
    ("Like", "Tap the heart on a profile that interests you.", CORAL),
    ("Mutual", "If they like you back, you're connected.", TEAL),
    ("Chat", "Friendly, low-pressure messaging opens up.", NAVY),
]
fx = Inches(5.6)
fy = Inches(2.40)
fw = Inches(7.3)
fh = Inches(1.30)
fgap = Inches(0.20)
for i, (head, body, accent) in enumerate(flow_steps):
    y = fy + (fh + fgap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, y, fw, fh)
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    box.shadow.inherit = False
    badge = s.shapes.add_shape(
        MSO_SHAPE.OVAL, fx + Inches(0.25), y + Inches(0.30),
        Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    badge.shadow.inherit = False
    add_text(s, fx + Inches(0.25), y + Inches(0.32), Inches(0.7),
             Inches(0.7), str(i + 1), size=24, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, fx + Inches(1.10), y + Inches(0.20), fw - Inches(1.3),
             Inches(0.5), head, size=22, bold=True, color=NAVY)
    add_text(s, fx + Inches(1.10), y + Inches(0.65), fw - Inches(1.3),
             Inches(0.55), body, size=15, color=INK)
add_text(s, fx, fy + (fh + fgap) * 3 + Inches(0.05), fw, Inches(0.35),
         "Block and report tools available at any time.",
         size=13, color=INK, align=PP_ALIGN.LEFT)
add_page_number(s, 8)

# Slide 9: Accessibility & trust (text-only — screenshots moved to slide 8)
s = prs.slides.add_slide(BLANK)
add_bg(s, SOFT_LAV)
add_accent_bars(s)
add_header(s, "Accessibility & trust",
           "Comfort, control, and confidence")
# Callout strip: how Silvers Fun is different from typical dating apps.
callout = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.85),
    Inches(12.2), Inches(0.6))
callout.adjustments[0] = 0.30
callout.fill.solid()
callout.fill.fore_color.rgb = WHITE
callout.line.color.rgb = CORAL
callout.line.width = Pt(2)
callout.shadow.inherit = False
add_text(s, Inches(0.55), Inches(1.93), Inches(12.2), Inches(0.45),
         "Companionship-first, not swipe-first dating — "
         "built around shared interests, meetups, and low-pressure connection.",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
trust_cards = [
    ("Built for friendship",
     "Shared interests and meetups, not dating swipes.", CORAL),
    ("Low-pressure pace",
     "Connect at your own speed; no urgency, no judgments.", TEAL),
    ("Mutual connection",
     "Chat opens only when both sides agree.", NAVY),
    ("Comfortable to use",
     "Large readable text, careful testing on real devices.", CORAL),
]
tc_w = Inches(2.95)
tc_h = Inches(3.6)
tc_gap = Inches(0.15)
tc_start = Inches(0.55)
tc_top = Inches(2.65)
for i, (head, body, accent) in enumerate(trust_cards):
    x = tc_start + (tc_w + tc_gap) * i
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, tc_top, tc_w, tc_h)
    card.adjustments[0] = 0.07
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(2)
    card.shadow.inherit = False
    stripe = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, tc_top, tc_w, Inches(0.25))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    stripe.shadow.inherit = False
    add_text(s, x + Inches(0.2), tc_top + Inches(0.5), tc_w - Inches(0.4),
             Inches(0.8), head, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), tc_top + Inches(1.4), tc_w - Inches(0.4),
             Inches(2.0), body, size=18, color=INK)
add_page_number(s, 9)

# Slide 10: Launch plan / CTA
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_accent_bars(s)
add_header(s, "Launch plan", "Listening before scaling")
phases = [
    ("Internal testing",
     "Team and trusted users on real devices.", CORAL),
    ("Closed testing",
     "Small groups of seniors and family members.", TEAL),
    ("Feedback loop",
     "Iterate on usability, safety, and warmth.", NAVY),
    ("Community outreach",
     "Senior centers, residences, and partners.", CORAL),
]
pw = Inches(2.95)
ph = Inches(2.6)
py = Inches(2.2)
for i, (h, b, c) in enumerate(phases):
    px = Inches(0.55) + (pw + Inches(0.15)) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = c
    box.line.width = Pt(2)
    box.shadow.inherit = False
    badge = s.shapes.add_shape(
        MSO_SHAPE.OVAL, px + Inches(0.25), py + Inches(0.25),
        Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = c
    badge.line.fill.background()
    badge.shadow.inherit = False
    add_text(s, px + Inches(0.25), py + Inches(0.27), Inches(0.55),
             Inches(0.55), str(i + 1), size=20, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, px + Inches(0.95), py + Inches(0.32), pw - Inches(1.1),
             Inches(0.6), h, size=18, bold=True, color=NAVY)
    add_text(s, px + Inches(0.25), py + Inches(1.1), pw - Inches(0.4),
             Inches(1.4), b, size=15, color=INK)

cta = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.2),
    Inches(12.2), Inches(1.6))
cta.adjustments[0] = 0.25
cta.fill.solid()
cta.fill.fore_color.rgb = NAVY
cta.line.fill.background()
cta.shadow.inherit = False
add_text(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(0.6),
         "Help us bring Silvers Fun to more seniors.",
         size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.6),
         "Partner with us, share with families, or join testing.",
         size=20, color=CREAM, align=PP_ALIGN.CENTER)
add_page_number(s, 10)

OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)

summary = """# Silvers Fun — Value Deck Summary

**File:** `docs/presentations/silvers_fun_value_deck.pptx`
**Format:** 16:9 widescreen, 10 editable slides

## Slides

1. **Title** — Silvers Fun: Helping seniors connect, meet, and enjoy life together. App icon and feature graphic.
2. **The problem** — Many seniors want companionship; existing apps don't fit them.
3. **Who it is for** — Seniors 65+, families, senior communities, hobby groups.
4. **What it does** — Value statement + 6-feature grid (Profiles, Discover, Meetups, Likes, Chat, English/Korean).
5. **Getting started** — One-tap Google Sign-In, profile setup, then discover and meet. Includes a placeholder for the sign-in screen.
6. **Explore people and activities** — Discover + Meetups + Create meetup screenshots.
7. **Personalize your experience** — Profile + Photos + Language screenshots, plus a placeholder for the Likes screen.
8. **Connect comfortably and safely** — Chat screenshot + Like → Mutual → Chat flow.
9. **Accessibility & trust** — Text-only. Callout: "Companionship-first, not swipe-first dating." Cards: built for friendship, low-pressure pace, mutual connection, comfortable to use.
10. **Launch plan / CTA** — Internal → closed testing → feedback → community outreach.

## Screenshots used

- `01_discover.jpg` (slide 6)
- `03_meetups.jpg` (slide 6)
- `04_create_meetup.jpg` (slide 6)
- `06_you_profile.jpg` (slide 7)
- `08_edit_photos.jpg` (slide 7)
- `07_language_picker.jpg` (slide 7)
- `05_chat.jpg` (slide 8)

## Screenshots to be added

- **Sign-in screen** (slide 5) — placeholder phone is in place; drop in the screenshot when available.
- **Likes screen** (slide 7) — placeholder phone is in place; drop in the screenshot when available.

Not used in the deck (kept available for future edits): `02_profile_detail.jpg`.

## Branding assets used

- `assets/branding/app_icon_play_store_512.png` (slide 1)
- `assets/branding/feature_graphic_1024x500.png` (slide 1)

## Style

- Backgrounds: cream (`#FAF4E8`) and soft lavender (`#EFEAF7`).
- Accents: navy (`#1F3A5F`), coral (`#E27D60`), teal (`#2A9D8F`).
- Large readable text, short bullets, rounded phone frames around screenshots.
- Positioned as a social companion and meetup app — not dating, no clinical claims.
"""

OUT_MD.write_text(summary, encoding="utf-8")
print(f"WROTE: {OUT_PPTX}")
print(f"WROTE: {OUT_MD}")
