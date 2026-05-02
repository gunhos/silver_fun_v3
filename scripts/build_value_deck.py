#!/usr/bin/env python3
"""Build the Silvers Fun value deck (PPTX).

Run from repo root:
    python3 scripts/build_value_deck.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw

# ---------- Paths ----------
OUT_DIR = os.path.join("docs", "presentations")
OUT_PPTX = os.path.join(OUT_DIR, "silvers_fun_value_deck.pptx")
SHOT_DIR = os.path.join("docs", "assets", "screenshots")
BRAND_DIR = os.path.join("assets", "branding")
TMP_DIR = "/tmp/silvers_fun_deck_assets"

os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(TMP_DIR, exist_ok=True)

# ---------- Brand palette ----------
LAVENDER       = RGBColor(0xEC, 0xE6, 0xF5)
LAVENDER_LIGHT = RGBColor(0xF6, 0xF2, 0xFB)
LAVENDER_DEEP  = RGBColor(0xDC, 0xD2, 0xEE)
CREAM          = RGBColor(0xFA, 0xF3, 0xE4)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
NAVY           = RGBColor(0x1F, 0x1B, 0x4D)
NAVY_SOFT      = RGBColor(0x3D, 0x3A, 0x6B)
SLATE          = RGBColor(0x55, 0x57, 0x77)
CORAL          = RGBColor(0xF3, 0x6F, 0x6F)
CORAL_SOFT     = RGBColor(0xFB, 0xD6, 0xD6)
TEAL           = RGBColor(0x1F, 0x6F, 0x75)
TEAL_SOFT      = RGBColor(0xC4, 0xDF, 0xE2)
SAGE           = RGBColor(0xA8, 0xCC, 0xC0)
YELLOW         = RGBColor(0xF4, 0xC7, 0x5A)
PHONE_BEZEL    = RGBColor(0x14, 0x14, 0x2C)

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

# ---------- Image preprocessing ----------

def round_corners(img, radius):
    img = img.convert("RGBA")
    w, h = img.size
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w, h), radius=radius, fill=255)
    out = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)
    return out

def prep_screenshot(filename, top_crop=110, bottom_crop=120):
    """Crop the Android status/gesture bars and round corners."""
    src = os.path.join(SHOT_DIR, filename)
    dst = os.path.join(TMP_DIR, filename.replace(".jpg", ".png"))
    img = Image.open(src)
    w, h = img.size
    img = img.crop((0, top_crop, w, h - bottom_crop))
    img = round_corners(img, int(min(img.size) * 0.06))
    img.save(dst, "PNG")
    return dst

SHOTS = {fn: prep_screenshot(fn) for fn in sorted(os.listdir(SHOT_DIR)) if fn.lower().endswith(".jpg")}
APP_ICON      = os.path.join(BRAND_DIR, "app_icon_play_store_512.png")
FEATURE_GRAPH = os.path.join(BRAND_DIR, "feature_graphic_1024x500.png")

# ---------- Presentation ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------

def _no_shadow(shape):
    shape.shadow.inherit = False

def add_bg(slide, color=LAVENDER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    _no_shadow(bg)
    return bg

def add_decor_band(slide, color=SAGE, height_in=0.45):
    """Subtle bottom band for warmth."""
    h = Inches(height_in)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - h, SW, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    _no_shadow(shape)
    return shape

def add_textbox(slide, x, y, w, h, text, *, font_size=18, color=NAVY, bold=False,
                font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_paragraphs(slide, x, y, w, h, items, *, font_size=18, color=NAVY,
                   font=FONT_BODY, align=PP_ALIGN.LEFT, line_spacing=1.25,
                   anchor=MSO_ANCHOR.TOP, bold=False):
    """Items: list of strings or (text, opts) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("font_size", font_size))
        run.font.bold = opts.get("bold", bold)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", color)
    return tb

def add_phone(slide, image_path, x_left, y_top, height_in):
    """Render a screenshot inside a rounded dark phone bezel."""
    img = Image.open(image_path)
    aspect = img.size[0] / img.size[1]
    bez_h = Inches(height_in)
    pad = Inches(0.10)
    img_h = bez_h - 2 * pad
    img_w = int(img_h * aspect)
    bez_w = img_w + 2 * pad
    bezel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_left, y_top, bez_w, bez_h)
    bezel.adjustments[0] = 0.07
    bezel.fill.solid()
    bezel.fill.fore_color.rgb = PHONE_BEZEL
    bezel.line.fill.background()
    _no_shadow(bezel)
    slide.shapes.add_picture(image_path, x_left + pad, y_top + pad, width=img_w, height=img_h)
    return bez_w, bez_h

def add_card(slide, x, y, w, h, *, fill=WHITE, border=None, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.25)
    else:
        shape.line.fill.background()
    _no_shadow(shape)
    return shape

def add_circle(slide, cx, cy, diameter_in, fill_color, *, text=None,
               text_color=WHITE, font_size=22, bold=True):
    d = Inches(diameter_in)
    x = int(cx - d / 2)
    y = int(cy - d / 2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _no_shadow(shape)
    if text is not None:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_TITLE
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape

def add_pill(slide, x, y, w, h, text, *, fill=WHITE, text_color=NAVY,
             font_size=16, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    _no_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape

def add_accent_line(slide, x, y, w, color=CORAL, height_in=0.05):
    h = Inches(height_in)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    _no_shadow(shape)
    return shape

def add_brandmark(slide, x=None, y=None, size_in=0.55):
    s = Inches(size_in)
    if x is None: x = SW - s - Inches(0.45)
    if y is None: y = Inches(0.4)
    return slide.shapes.add_picture(APP_ICON, x, y, width=s, height=s)

def add_slide_title(slide, title, subtitle=None, *, y=Inches(0.5)):
    add_textbox(slide, Inches(0.7), y, Inches(11.0), Inches(0.85), title,
                font_size=38, color=NAVY, bold=True, font=FONT_TITLE)
    if subtitle:
        add_textbox(slide, Inches(0.7), y + Inches(0.85), Inches(11.0), Inches(0.5),
                    subtitle, font_size=18, color=SLATE, italic=True)
        accent_y = y + Inches(0.85 + 0.55)
    else:
        accent_y = y + Inches(0.85)
    add_accent_line(slide, Inches(0.7), accent_y, Inches(0.9), color=CORAL)

def add_footer(slide, page_num):
    add_textbox(slide, Inches(0.45), SH - Inches(0.42), Inches(8), Inches(0.3),
                "Silvers Fun  ·  Friendly company for the next chapter",
                font_size=10, color=SLATE)
    add_textbox(slide, SW - Inches(1.2), SH - Inches(0.42), Inches(0.8), Inches(0.3),
                f"{page_num} / 10", font_size=10, color=SLATE, align=PP_ALIGN.RIGHT)

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ---------- Slide 1: Title ----------

def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LAVENDER)
    # decorative sage band at the bottom (matches brand graphic)
    add_decor_band(s, SAGE, height_in=0.55)
    # warm sun accent in upper right
    sun = s.shapes.add_shape(MSO_SHAPE.OVAL, SW - Inches(2.2), Inches(0.6),
                             Inches(1.1), Inches(1.1))
    sun.fill.solid(); sun.fill.fore_color.rgb = YELLOW
    sun.line.fill.background(); _no_shadow(sun)

    # Big app icon
    icon_w = Inches(2.2)
    s.shapes.add_picture(APP_ICON, (SW - icon_w) / 2, Inches(0.95),
                         width=icon_w, height=icon_w)
    # Title
    add_textbox(s, Inches(0), Inches(3.45), SW, Inches(1.2),
                "Silvers Fun", font_size=80, color=NAVY, bold=True,
                font=FONT_TITLE, align=PP_ALIGN.CENTER)
    # Tagline
    add_textbox(s, Inches(0), Inches(4.55), SW, Inches(0.7),
                "Friendly company for the next chapter",
                font_size=28, color=CORAL, italic=True,
                align=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(s, Inches(0), Inches(5.35), SW, Inches(0.5),
                "A warm social companion app for adults 65+",
                font_size=18, color=SLATE, align=PP_ALIGN.CENTER)
    # Pills row
    pills = ["Connect", "Meetups", "Enjoy life"]
    pill_w, pill_h = Inches(1.7), Inches(0.5)
    gap = Inches(0.3)
    total = pill_w * len(pills) + gap * (len(pills) - 1)
    px = (SW - total) / 2
    for i, label in enumerate(pills):
        accent = TEAL if i % 2 == 0 else CORAL
        add_pill(s, px + i * (pill_w + gap), Inches(6.05), pill_w, pill_h,
                 label, fill=WHITE, text_color=accent, font_size=15)
    set_notes(s,
        "Welcome and intro. Silvers Fun is a warm social companion app built for adults 65+. "
        "The tagline 'Friendly company for the next chapter' captures the goal: help seniors find "
        "real friendship, shared interests, and group meetups in a calm, low-pressure space. "
        "This deck explains the problem, the people we serve, the product, and how partners and "
        "communities can help us launch.")

# ---------- Slide 2: The Problem ----------

def slide_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "Later in life can be lonely",
                    subtitle="The world doesn't stop changing — but the easy ways to meet new people often do.")

    # Three empathy cards
    card_w, card_h = Inches(3.7), Inches(3.4)
    gap = Inches(0.35)
    total = card_w * 3 + gap * 2
    start_x = (SW - total) / 2
    y = Inches(2.55)
    cards = [
        ("Friends move", "People retire, relocate, or slow down. The old social rhythms quietly fade.", CORAL),
        ("Family gets busy", "Children and grandchildren lead full lives. Visits get rarer than anyone wants.", TEAL),
        ("Hobbies need company", "Walks, lunches, music, travel — the joys are bigger when shared with someone.", NAVY_SOFT),
    ]
    for i, (title, body, accent) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(s, x, y, card_w, card_h, fill=WHITE)
        add_circle(s, x + Inches(0.6), y + Inches(0.6), 0.55, accent,
                   text=str(i + 1), font_size=20)
        add_textbox(s, x + Inches(1.3), y + Inches(0.35), card_w - Inches(1.5), Inches(0.7),
                    title, font_size=22, color=NAVY, bold=True)
        add_textbox(s, x + Inches(0.45), y + Inches(1.4), card_w - Inches(0.9), card_h - Inches(1.6),
                    body, font_size=16, color=SLATE, line_spacing=1.3)
    # Closing line
    add_textbox(s, Inches(0.7), Inches(6.3), SW - Inches(1.4), Inches(0.5),
                "Most apps weren't built for this stage of life — or this kind of need.",
                font_size=18, color=NAVY, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 2)
    set_notes(s,
        "Loneliness in older adulthood is widely recognized but underserved. Friends move "
        "or pass on, family is busy, and hobbies often want company. Existing social apps "
        "are designed for younger users, fast scrolling, and dating — none of which fit "
        "the calm, friendship-first need we hear from seniors and their families.")

# ---------- Slide 3: Target Users ----------

def slide_audience():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "Made for adults 65+",
                    subtitle="Different stories, the same wish: warm, easy, friendly company.")

    # 3 persona cards
    card_w, card_h = Inches(3.7), Inches(3.5)
    gap = Inches(0.35)
    total = card_w * 3 + gap * 2
    start_x = (SW - total) / 2
    y = Inches(2.55)
    personas = [
        ("Newly retired", "Eager for a fresh weekly rhythm — coffee, walks, classes, new friends.", TEAL),
        ("Recently widowed", "Wants warm conversation and gentle company, with no pressure.", CORAL),
        ("Empty nesters", "Time to invest in friendships, hobbies, and groups they actually enjoy.", NAVY_SOFT),
    ]
    for i, (title, body, accent) in enumerate(personas):
        x = start_x + i * (card_w + gap)
        add_card(s, x, y, card_w, card_h, fill=WHITE)
        # accent strip on top
        add_card(s, x, y, card_w, Inches(0.18), fill=accent, radius=0.12)
        # heart-style badge
        add_circle(s, x + card_w / 2, y + Inches(1.0), 0.85, accent,
                   text="♥", font_size=30)
        add_textbox(s, x + Inches(0.3), y + Inches(1.65), card_w - Inches(0.6), Inches(0.6),
                    title, font_size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.4), y + Inches(2.3), card_w - Inches(0.8), card_h - Inches(2.5),
                    body, font_size=15, color=SLATE, line_spacing=1.35,
                    align=PP_ALIGN.CENTER)
    # Closing pill
    pill_w = Inches(7.5)
    add_pill(s, (SW - pill_w) / 2, Inches(6.45), pill_w, Inches(0.5),
             "All looking for friendly, low-pressure connection",
             fill=LAVENDER_DEEP, text_color=NAVY, font_size=15)
    add_footer(s, 3)
    set_notes(s,
        "Our core audience is adults 65+. Three common situations bring people in: newly "
        "retired, recently widowed, and empty nesters. The unifying need is friendly, "
        "low-pressure company — not romance, not networking, not entertainment. We design "
        "every interaction around that.")

# ---------- Slide 4: Core Value Proposition ----------

def slide_value():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LAVENDER_LIGHT)
    add_brandmark(s)
    # Title
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(11.0), Inches(0.9),
                "A warm place built for the way seniors want to socialize",
                font_size=34, color=NAVY, bold=True, font=FONT_TITLE)
    add_accent_line(s, Inches(0.7), Inches(1.45), Inches(0.9), color=CORAL)

    # Big value statement
    add_textbox(s, Inches(0.9), Inches(1.85), SW - Inches(1.8), Inches(0.9),
                '"Find people you actually like. Do things you actually enjoy. Together."',
                font_size=24, color=NAVY_SOFT, italic=True, align=PP_ALIGN.CENTER)

    # Three pillar cards
    pillars = [
        ("Connect", "Discover people whose interests, neighborhood, and pace fit yours.", TEAL, "♥"),
        ("Meet up", "Join or host casual group activities — coffee, walks, music, more.", CORAL, "☕"),
        ("Enjoy life", "Gentle, calm experience with large readable screens and clear actions.", NAVY_SOFT, "✿"),
    ]
    card_w, card_h = Inches(3.8), Inches(2.9)
    gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start_x = (SW - total) / 2
    y = Inches(3.05)
    for i, (title, body, accent, glyph) in enumerate(pillars):
        x = start_x + i * (card_w + gap)
        add_card(s, x, y, card_w, card_h, fill=WHITE)
        add_circle(s, x + card_w / 2, y + Inches(0.7), 1.0, accent,
                   text=glyph, font_size=34)
        add_textbox(s, x + Inches(0.3), y + Inches(1.45), card_w - Inches(0.6), Inches(0.55),
                    title, font_size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.4), y + Inches(2.05), card_w - Inches(0.8), Inches(0.85),
                    body, font_size=14, color=SLATE, line_spacing=1.3, align=PP_ALIGN.CENTER)

    # Footer band with feature graphic-inspired accents
    add_textbox(s, Inches(0.7), Inches(6.45), SW - Inches(1.4), Inches(0.5),
                "Friendship-first, interest-led, and built for real life off the screen.",
                font_size=18, color=NAVY, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 4)
    set_notes(s,
        "The promise in one sentence: find people you like, do things you enjoy, together. "
        "Three pillars guide the product — Connect (discover compatible people), Meet up "
        "(turn connections into real activities), Enjoy life (a calm, accessible experience). "
        "Everything else flows from these three.")

# ---------- Slide 5: Main Features ----------

def slide_features():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "Everything you need to meet good people",
                    subtitle="Six simple capabilities, designed for clarity over cleverness.")

    # Two phones on the right (profile + language picker), side by side
    phone_h = 4.0
    sample = Image.open(SHOTS["02_profile_detail.jpg"])
    aspect = sample.size[0] / sample.size[1]
    img_h = Inches(phone_h) - Inches(0.20)
    img_w = int(img_h * aspect)
    bez_w = img_w + Inches(0.20)
    phone_gap = Inches(0.25)
    phones_total = bez_w * 2 + phone_gap
    phones_x = SW - phones_total - Inches(0.5)
    phones_y = Inches(2.4)
    add_phone(s, SHOTS["02_profile_detail.jpg"], phones_x, phones_y, phone_h)
    add_phone(s, SHOTS["07_language_picker.jpg"], phones_x + bez_w + phone_gap, phones_y, phone_h)
    cap_y = phones_y + Inches(phone_h) + Inches(0.08)
    add_textbox(s, phones_x - Inches(0.2), cap_y, bez_w + Inches(0.4), Inches(0.3),
                "Profile", font_size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, phones_x + bez_w + phone_gap - Inches(0.2), cap_y,
                bez_w + Inches(0.4), Inches(0.3),
                "Languages", font_size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)

    # Features grid 2x3 on the left
    feats = [
        ("Profiles", "Photos, age, and personal interests.", TEAL),
        ("Discover", "Browse people who fit your interests.", CORAL),
        ("Meetups", "Create or join casual group activities.", NAVY_SOFT),
        ("Likes & connections", "Mutual likes open a friendly chat.", TEAL),
        ("Chat", "Simple in-app messaging once connected.", CORAL),
        ("English & 한국어", "Full support for both languages.", NAVY_SOFT),
    ]
    grid_x = Inches(0.7)
    grid_y = Inches(2.4)
    avail_w = phones_x - grid_x - Inches(0.3)
    col_gap = Inches(0.2)
    cell_w = (avail_w - col_gap) / 2
    cell_h = Inches(1.15)
    row_gap = Inches(0.18)
    for i, (title, body, accent) in enumerate(feats):
        col = i % 2
        row = i // 2
        x = grid_x + col * (cell_w + col_gap)
        y = grid_y + row * (cell_h + row_gap)
        add_card(s, x, y, cell_w, cell_h, fill=WHITE)
        add_circle(s, x + Inches(0.5), y + cell_h / 2, 0.55, accent,
                   text=str(i + 1), font_size=16)
        add_textbox(s, x + Inches(0.95), y + Inches(0.15), cell_w - Inches(1.1), Inches(0.4),
                    title, font_size=15, color=NAVY, bold=True)
        add_textbox(s, x + Inches(0.95), y + Inches(0.55), cell_w - Inches(1.1), Inches(0.55),
                    body, font_size=11, color=SLATE, line_spacing=1.25)
    add_footer(s, 5)
    set_notes(s,
        "Six capabilities cover the full experience: profiles with interests, a Discover "
        "feed of compatible people, group Meetups, mutual Likes & Connections, in-app Chat, "
        "and full Korean/English support. Each one is intentionally simple — clarity over "
        "cleverness. Two example screens on the right: a profile (interests-first, not "
        "swipe-first) and the language picker showing built-in English and Korean.")

# ---------- Slide 6: Product walkthrough ----------

def slide_walkthrough():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "A quick look inside",
                    subtitle="Discover people, then browse meetups — all in a calm, readable layout.")

    # Two phones side by side
    phone_h = 4.0
    sample = Image.open(SHOTS["01_discover.jpg"])
    aspect = sample.size[0] / sample.size[1]
    img_h = Inches(phone_h) - Inches(0.20)
    img_w = int(img_h * aspect)
    bez_w = img_w + Inches(0.20)
    gap = Inches(2.0)
    total = bez_w * 2 + gap
    start_x = (SW - total) / 2
    y = Inches(2.45)
    add_phone(s, SHOTS["01_discover.jpg"], start_x, y, phone_h)
    add_phone(s, SHOTS["03_meetups.jpg"], start_x + bez_w + gap, y, phone_h)

    # Captions
    cap_y = y + Inches(phone_h) + Inches(0.1)
    add_textbox(s, start_x - Inches(0.5), cap_y, bez_w + Inches(1.0), Inches(0.32),
                "Discover",
                font_size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, start_x - Inches(0.5), cap_y + Inches(0.34),
                bez_w + Inches(1.0), Inches(0.28),
                "Browse profiles, tap a heart to like.",
                font_size=12, color=SLATE, align=PP_ALIGN.CENTER, italic=True)
    add_textbox(s, start_x + bez_w + gap - Inches(0.5), cap_y,
                bez_w + Inches(1.0), Inches(0.32),
                "Meetups",
                font_size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, start_x + bez_w + gap - Inches(0.5), cap_y + Inches(0.34),
                bez_w + Inches(1.0), Inches(0.28),
                "See group activities you can join.",
                font_size=12, color=SLATE, align=PP_ALIGN.CENTER, italic=True)

    # Connector arrow between phones
    arrow_y = y + Inches(phone_h) / 2
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               start_x + bez_w + Inches(0.2), arrow_y - Inches(0.25),
                               gap - Inches(0.4), Inches(0.5))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = CORAL_SOFT
    arrow.line.fill.background(); _no_shadow(arrow)
    add_footer(s, 6)
    set_notes(s,
        "Two of the main screens. Discover shows one profile at a time with a clear like "
        "action — no rapid swiping. Meetups is a feed of group activities people can join "
        "or create. The bottom navigation (Discover, Liked you, Meetups, Chats, You) is "
        "constant so users always know where they are.")

# ---------- Slide 7: User journey ----------

def slide_journey():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "From hello to a real meetup",
                    subtitle="A simple four-step path — every step optional, none of it pushy.")

    # 4 step badges horizontally — left side, taking ~8.5" width
    steps = [
        ("Sign up", "Add photos & the interests that matter to you."),
        ("Discover", "Browse people who share your interests."),
        ("Connect", "Mutual likes turn into a friendly chat."),
        ("Meet up", "Join or create a group activity."),
    ]

    # Phones occupy the right side (steps 3 & 4 visualized)
    phone_h = 4.0
    sample = Image.open(SHOTS["05_chat.jpg"])
    aspect = sample.size[0] / sample.size[1]
    img_h = Inches(phone_h) - Inches(0.20)
    img_w = int(img_h * aspect)
    bez_w = img_w + Inches(0.20)
    p_gap = Inches(0.25)
    phones_total = bez_w * 2 + p_gap
    px = SW - phones_total - Inches(0.5)
    py = Inches(2.45)
    add_phone(s, SHOTS["05_chat.jpg"], px, py, phone_h)
    add_phone(s, SHOTS["04_create_meetup.jpg"], px + bez_w + p_gap, py, phone_h)
    add_textbox(s, px - Inches(0.2), py + Inches(phone_h) + Inches(0.08),
                bez_w + Inches(0.4), Inches(0.3),
                "Connect & chat", font_size=12, color=SLATE,
                italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, px + bez_w + p_gap - Inches(0.2),
                py + Inches(phone_h) + Inches(0.08),
                bez_w + Inches(0.4), Inches(0.3),
                "Create a meetup", font_size=12, color=SLATE,
                italic=True, align=PP_ALIGN.CENTER)

    # Steps stacked vertically on the left
    steps_x = Inches(0.7)
    avail_w = px - steps_x - Inches(0.4)
    steps_y = Inches(2.45)
    step_h = Inches(0.95)
    step_gap = Inches(0.18)
    accents = [TEAL, CORAL, TEAL, CORAL]
    for i, (title, body) in enumerate(steps):
        cy = steps_y + i * (step_h + step_gap)
        add_circle(s, steps_x + Inches(0.45), cy + step_h / 2, 0.75,
                   accents[i], text=str(i + 1), font_size=22)
        add_textbox(s, steps_x + Inches(1.05), cy + Inches(0.05),
                    avail_w - Inches(1.05), Inches(0.45),
                    title, font_size=20, color=NAVY, bold=True)
        add_textbox(s, steps_x + Inches(1.05), cy + Inches(0.5),
                    avail_w - Inches(1.05), Inches(0.45),
                    body, font_size=13, color=SLATE, line_spacing=1.25)
        # vertical connector line
        if i < 3:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      steps_x + Inches(0.42),
                                      cy + step_h - Inches(0.05),
                                      Inches(0.06),
                                      step_gap + Inches(0.1))
            line.fill.solid(); line.fill.fore_color.rgb = LAVENDER_DEEP
            line.line.fill.background(); _no_shadow(line)

    add_footer(s, 7)
    set_notes(s,
        "The journey is intentionally short. Sign up and pick interests, browse Discover, "
        "exchange mutual likes to start a chat, then create or join a real meetup. There is "
        "no pressure to rush any step — many users will browse for a while before "
        "connecting, which is fine.")

# ---------- Slide 8: Different from dating apps ----------

def slide_not_dating():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "Not a dating app — a companion app",
                    subtitle="We built every interaction around friendship, calm, and accessibility.")

    # Two-column comparison cards
    col_w = Inches(5.6)
    col_h = Inches(4.2)
    gap = Inches(0.4)
    total = col_w * 2 + gap
    start_x = (SW - total) / 2
    y = Inches(2.4)

    # Left: typical dating apps
    add_card(s, start_x, y, col_w, col_h, fill=WHITE, border=LAVENDER_DEEP)
    add_textbox(s, start_x + Inches(0.4), y + Inches(0.3), col_w - Inches(0.8), Inches(0.5),
                "Typical dating apps", font_size=22, color=SLATE, bold=True)
    add_paragraphs(s, start_x + Inches(0.4), y + Inches(1.0),
                   col_w - Inches(0.8), col_h - Inches(1.3),
                   ["•  Fast swiping, snap judgements",
                    "•  Endless options, decision fatigue",
                    "•  Romance-focused vocabulary",
                    "•  Designed for users in their 20s–40s",
                    "•  Small text, dense screens"],
                   font_size=16, color=SLATE, line_spacing=1.55)

    # Right: Silvers Fun
    add_card(s, start_x + col_w + gap, y, col_w, col_h, fill=CREAM, border=CORAL)
    add_textbox(s, start_x + col_w + gap + Inches(0.4), y + Inches(0.3),
                col_w - Inches(0.8), Inches(0.5),
                "Silvers Fun", font_size=22, color=CORAL, bold=True)
    add_paragraphs(s, start_x + col_w + gap + Inches(0.4), y + Inches(1.0),
                   col_w - Inches(0.8), col_h - Inches(1.3),
                   ["•  Calm browsing, one profile at a time",
                    "•  Curated by interests and meetups",
                    "•  Friendship and shared activity vocabulary",
                    "•  Designed specifically for adults 65+",
                    "•  Large readable type, generous spacing"],
                   font_size=16, color=NAVY, line_spacing=1.55)

    # Closing line
    add_textbox(s, Inches(0.7), Inches(6.7), SW - Inches(1.4), Inches(0.35),
                "We're about company, not chemistry.",
                font_size=20, color=NAVY, italic=True, bold=True,
                align=PP_ALIGN.CENTER)
    add_footer(s, 8)
    set_notes(s,
        "We get asked often: 'Is this a dating app?' No. Dating apps optimise for fast "
        "judgement and romance. Silvers Fun optimises for warm friendship, shared interests, "
        "and real-world group activities. The screen layouts, copy, pace, and accessibility "
        "all follow from that. One sentence to remember: we're about company, not chemistry.")

# ---------- Slide 9: Trust and safety ----------

def slide_trust():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_brandmark(s)
    add_slide_title(s, "Built with care and respect",
                    subtitle="Designed for control, calm, and clarity.")

    # Two phones on the right
    phone_h = 4.1
    sample = Image.open(SHOTS["08_edit_photos.jpg"])
    aspect = sample.size[0] / sample.size[1]
    img_h = Inches(phone_h) - Inches(0.20)
    img_w = int(img_h * aspect)
    bez_w = img_w + Inches(0.20)
    gap = Inches(0.3)
    phones_total = bez_w * 2 + gap
    px = SW - phones_total - Inches(0.6)
    py = Inches(2.4)
    add_phone(s, SHOTS["06_you_profile.jpg"], px, py, phone_h)
    add_phone(s, SHOTS["08_edit_photos.jpg"], px + bez_w + gap, py, phone_h)
    add_textbox(s, px - Inches(0.2), py + Inches(phone_h) + Inches(0.08),
                bez_w + Inches(0.4), Inches(0.28),
                "Manage your profile", font_size=11, color=SLATE,
                italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, px + bez_w + gap - Inches(0.2),
                py + Inches(phone_h) + Inches(0.08),
                bez_w + Inches(0.4), Inches(0.28),
                "Edit your photos", font_size=11, color=SLATE,
                italic=True, align=PP_ALIGN.CENTER)

    # Safety bullets on the left
    items_x = Inches(0.7)
    items_y = Inches(2.4)
    items_w = px - items_x - Inches(0.4)

    bullets = [
        ("Photo control", "Choose what to share, edit or remove anytime."),
        ("Mutual likes only", "Chat opens after both people show interest — no unwanted DMs."),
        ("Private by default", "Personal contact details stay off-platform unless you share them."),
        ("Easy profile management", "Update interests, language, and photos from one place."),
    ]
    item_h = Inches(0.85)
    for i, (head, body) in enumerate(bullets):
        cy = items_y + i * (item_h + Inches(0.16))
        add_circle(s, items_x + Inches(0.35), cy + item_h / 2, 0.55,
                   TEAL if i % 2 == 0 else CORAL, text="✓", font_size=18)
        add_textbox(s, items_x + Inches(0.85), cy + Inches(0.02),
                    items_w - Inches(0.85), Inches(0.4),
                    head, font_size=17, color=NAVY, bold=True)
        add_textbox(s, items_x + Inches(0.85), cy + Inches(0.45),
                    items_w - Inches(0.85), Inches(0.4),
                    body, font_size=13, color=SLATE, line_spacing=1.25)

    # Disclaimer line
    add_textbox(s, Inches(0.7), Inches(6.7), SW - Inches(1.4), Inches(0.3),
                "Trust and safety practices evolve as we learn from testers and partners.",
                font_size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 9)
    set_notes(s,
        "Trust matters more here than in most apps. Today's design includes user-controlled "
        "photos, mutual-like-required chat, private personal details by default, and easy "
        "profile management. We expect this list to grow as we learn from testers and senior "
        "communities — moderation policies, reporting flows, and family-side tools are all "
        "areas we want partner input on.")

# ---------- Slide 10: Launch & CTA ----------

def slide_launch():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LAVENDER_LIGHT)
    add_decor_band(s, SAGE, height_in=0.4)

    # Title
    add_textbox(s, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.85),
                "Help us bring this to life",
                font_size=38, color=NAVY, bold=True, font=FONT_TITLE)
    add_textbox(s, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.5),
                "Now in early testing. We're inviting partners, communities, and families to join.",
                font_size=18, color=SLATE, italic=True)
    add_accent_line(s, Inches(0.7), Inches(1.95), Inches(0.9), color=CORAL)

    # Phase timeline (3 cards across full width)
    phases = [
        ("Now", "Closed testing with a small group", TEAL),
        ("Next", "Beta with senior communities & partners", CORAL),
        ("Soon", "Public launch in English & 한국어", NAVY_SOFT),
    ]
    pcard_w = Inches(3.95)
    pgap = Inches(0.18)
    pstart = Inches(0.7)
    py = Inches(2.3)
    pheight = Inches(1.4)
    for i, (label, body, accent) in enumerate(phases):
        x = pstart + i * (pcard_w + pgap)
        add_card(s, x, py, pcard_w, pheight, fill=WHITE)
        add_card(s, x, py, Inches(0.18), pheight, fill=accent, radius=0.3)
        add_textbox(s, x + Inches(0.45), py + Inches(0.18), pcard_w - Inches(0.6), Inches(0.5),
                    label, font_size=20, color=accent, bold=True)
        add_textbox(s, x + Inches(0.45), py + Inches(0.7), pcard_w - Inches(0.6), Inches(0.6),
                    body, font_size=15, color=NAVY, line_spacing=1.3)

    # CTA box
    cta_x = Inches(0.7)
    cta_y = Inches(3.95)
    cta_w = Inches(11.95)
    cta_h = Inches(2.0)
    add_card(s, cta_x, cta_y, cta_w, cta_h, fill=CREAM, border=CORAL)
    add_textbox(s, cta_x + Inches(0.4), cta_y + Inches(0.2),
                cta_w - Inches(0.8), Inches(0.5),
                "How you can help", font_size=22, color=NAVY, bold=True)
    cta_items = [
        ("Become a tester", "If you're 65+ or care for someone who is, try the app and tell us what works."),
        ("Partner with us", "Senior centers, community groups, and family services — let's pilot together."),
        ("Share with someone", "Send Silvers Fun to a parent, neighbor, or friend who'd love a warmer way to connect."),
    ]
    item_w = (cta_w - Inches(0.8)) / 3
    iy = cta_y + Inches(0.8)
    for i, (head, body) in enumerate(cta_items):
        ix = cta_x + Inches(0.4) + i * item_w
        accent = [TEAL, CORAL, NAVY_SOFT][i]
        add_circle(s, ix + Inches(0.4), iy + Inches(0.25), 0.5, accent,
                   text=str(i + 1), font_size=16)
        add_textbox(s, ix + Inches(0.85), iy, item_w - Inches(1.0), Inches(0.4),
                    head, font_size=15, color=NAVY, bold=True)
        add_textbox(s, ix + Inches(0.85), iy + Inches(0.4),
                    item_w - Inches(1.0), Inches(0.7),
                    body, font_size=12, color=SLATE, line_spacing=1.3)

    # Closing row: app icon + wordmark
    icon_size = Inches(0.85)
    icon_x = Inches(5.3)
    closing_y = Inches(6.2)
    s.shapes.add_picture(APP_ICON, icon_x, closing_y, width=icon_size, height=icon_size)
    add_textbox(s, icon_x + icon_size + Inches(0.2), closing_y + Inches(0.18),
                Inches(3.5), Inches(0.5),
                "Silvers Fun", font_size=24, color=NAVY, bold=True)
    add_textbox(s, Inches(0.7), Inches(7.05), SW - Inches(1.4), Inches(0.3),
                "hello@silversfun.app",
                font_size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 10)
    set_notes(s,
        "Three phases: closed testing today, partner beta next, public launch in English "
        "and Korean soon. Three asks for the audience: become a tester, become a pilot "
        "partner, or share with a senior in your life. The contact line is illustrative — "
        "replace with your real contact before sharing externally.")

# ---------- Build ----------

slide_title()
slide_problem()
slide_audience()
slide_value()
slide_features()
slide_walkthrough()
slide_journey()
slide_not_dating()
slide_trust()
slide_launch()

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}")
print(f"Slides: {len(prs.slides)}")
